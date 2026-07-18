"""Markdown-like content -> PowerPoint (.pptx) renderer."""

from __future__ import annotations

import re
from pathlib import Path

_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)$")
_BULLET_RE = re.compile(r"^(\s*)[-*+]\s+(.*)$")
_NUMBERED_RE = re.compile(r"^\s*\d+[.)]\s+(.*)$")
_LINK_RE = re.compile(r"\[([^\]]+)\]\(([^)]+)\)")
_RULE_CHARS = {"-", "_", "*"}

# Slide layout indexes in the default python-pptx template.
_TITLE_AND_CONTENT_LAYOUT = 1
_TITLE_ONLY_LAYOUT = 5

MAX_BULLETS_PER_SLIDE = 10
BODY_FONT_PT = 18


def render_pptx_file(content: str, target: Path, *, title: str = "") -> int:
    """Render Markdown-like text into a .pptx file and return the bytes written.

    Each `#`/`##` heading starts a new slide titled by the heading; the lines
    below it become the slide's bullets. Content before the first heading (or
    heading-less content) lands on a slide titled with the fallback title.
    """
    presentation = _new_presentation()
    for slide in _split_slides(content or "", default_title=title or "Generated Presentation"):
        _add_content_slides(presentation, slide["title"], slide["lines"])
    presentation.save(str(target))
    return target.stat().st_size


def _new_presentation():
    try:
        from pptx import Presentation
    except ImportError as err:
        raise ValueError(
            "PowerPoint output requires the python-pptx package. Install it with: pip install python-pptx"
        ) from err
    return Presentation()


def _split_slides(content: str, *, default_title: str) -> list[dict]:
    slides: list[dict] = []
    current: dict = {"title": default_title, "lines": []}
    has_heading = False
    in_code_block = False

    for raw_line in content.splitlines():
        stripped = raw_line.strip()

        if stripped.startswith("```"):
            in_code_block = not in_code_block
            continue
        if not stripped:
            continue
        if in_code_block:
            current["lines"].append((1, stripped))
            continue

        heading = _HEADING_RE.match(stripped)
        if heading and len(heading.group(1)) <= 2:
            if has_heading or current["lines"]:
                slides.append(current)
            current = {"title": _plain_text(heading.group(2)), "lines": []}
            has_heading = True
            continue
        if heading:
            current["lines"].append((0, _plain_text(heading.group(2))))
            continue
        if len(stripped) >= 3 and set(stripped) <= _RULE_CHARS:
            continue
        bullet = _BULLET_RE.match(raw_line)
        if bullet:
            level = 1 if len(bullet.group(1)) >= 2 else 0
            current["lines"].append((level, _plain_text(bullet.group(2))))
            continue
        numbered = _NUMBERED_RE.match(raw_line)
        if numbered:
            current["lines"].append((0, _plain_text(numbered.group(1))))
            continue
        current["lines"].append((0, _plain_text(stripped)))

    slides.append(current)
    return slides


def _add_content_slides(presentation, title: str, lines: list[tuple[int, str]]) -> None:
    chunks = [lines[index : index + MAX_BULLETS_PER_SLIDE] for index in range(0, len(lines), MAX_BULLETS_PER_SLIDE)]
    for index, chunk in enumerate(chunks or [[]]):
        _add_slide(presentation, title if index == 0 else f"{title} (cont.)", chunk)


def _add_slide(presentation, title: str, lines: list[tuple[int, str]]) -> None:
    from pptx.util import Pt

    layout_index = _TITLE_AND_CONTENT_LAYOUT if lines else _TITLE_ONLY_LAYOUT
    slide = presentation.slides.add_slide(presentation.slide_layouts[layout_index])
    slide.shapes.title.text = title
    if not lines:
        return

    body = slide.placeholders[1].text_frame
    body.word_wrap = True
    for index, (level, text) in enumerate(lines):
        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
        paragraph.text = text
        paragraph.level = min(level, 4)
        for run in paragraph.runs:
            run.font.size = Pt(BODY_FONT_PT)


def _plain_text(text: str) -> str:
    return _LINK_RE.sub(r"\1 (\2)", text).replace("**", "").replace("`", "").strip()
