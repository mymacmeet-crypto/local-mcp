import unittest
import uuid
from pathlib import Path
from unittest.mock import patch

from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from local_mcp.file_generation import append_generated_file, write_generated_file
from local_mcp.tools.file_generation import generate_file


class FileGenerationTests(unittest.TestCase):
    def tempdir(self):
        root = Path.cwd() / ".tmp" / "unit-tests" / f"{self._testMethodName}-{uuid.uuid4().hex}"
        root.mkdir(parents=True)
        return str(root)

    def test_write_generated_file_creates_markdown_under_env_output_dir(self):
        tmp = self.tempdir()
        with patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True):
            result = write_generated_file("notes/report", "# Title")

        expected = Path(tmp).resolve() / "notes" / "report.md"
        self.assertEqual(result.path, expected)
        self.assertEqual(expected.read_text(encoding="utf-8"), "# Title\n")
        self.assertEqual(result.file_type, "md")
        self.assertFalse(result.overwritten)

    def test_write_generated_file_refuses_existing_file_without_overwrite(self):
        tmp = self.tempdir()
        target = Path(tmp) / "notes.md"
        target.write_text("old", encoding="utf-8")

        with patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True):
            with self.assertRaises(ValueError):
                write_generated_file("notes.md", "new")

            result = write_generated_file("notes.md", "new", overwrite=True)
        self.assertTrue(result.overwritten)
        self.assertEqual(target.read_text(encoding="utf-8"), "new\n")

    def test_write_generated_file_rejects_path_traversal(self):
        tmp = self.tempdir()
        with patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True):
            with self.assertRaises(ValueError):
                write_generated_file("../escape", "content")

    def test_write_generated_file_rejects_mismatched_suffix(self):
        tmp = self.tempdir()
        with patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True):
            with self.assertRaises(ValueError):
                write_generated_file("notes.txt", "content", file_type="pdf")

    def test_write_generated_file_rejects_unsupported_file_type(self):
        tmp = self.tempdir()
        with patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True):
            with self.assertRaises(ValueError):
                write_generated_file("notes", "content", file_type="xlsx")

    def test_write_generated_file_creates_plain_text_file(self):
        tmp = self.tempdir()
        with patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True):
            result = write_generated_file("notes/plain", "plain text body", file_type="txt")

        expected = Path(tmp).resolve() / "notes" / "plain.txt"
        self.assertEqual(result.path, expected)
        self.assertEqual(result.file_type, "txt")
        self.assertEqual(expected.read_text(encoding="utf-8"), "plain text body\n")

    def test_write_generated_file_creates_pdf_under_env_output_dir(self):
        tmp = self.tempdir()
        with patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True):
            result = write_generated_file(
                "reports/summary",
                "# Summary\n\n- First point\n- Second point",
                file_type="pdf",
            )

        expected = Path(tmp).resolve() / "reports" / "summary.pdf"
        self.assertEqual(result.path, expected)
        self.assertEqual(result.file_type, "pdf")
        self.assertGreater(result.bytes_written, 0)
        text = "\n".join(page.extract_text() or "" for page in PdfReader(str(expected)).pages)
        self.assertIn("Summary", text)
        self.assertIn("First point", text)

    def test_write_generated_file_creates_word_document(self):
        tmp = self.tempdir()
        with patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True):
            result = write_generated_file(
                "reports/word-report",
                "# Word Report\n\nIntro paragraph.\n\n- First bullet\n- Second bullet",
                file_type="docx",
            )

        expected = Path(tmp).resolve() / "reports" / "word-report.docx"
        self.assertEqual(result.path, expected)
        self.assertEqual(result.file_type, "docx")
        self.assertGreater(result.bytes_written, 0)
        texts = [paragraph.text for paragraph in Document(str(expected)).paragraphs]
        self.assertIn("Word Report", texts)
        self.assertIn("Intro paragraph.", texts)
        self.assertIn("First bullet", texts)

    def test_write_generated_file_creates_powerpoint_presentation(self):
        tmp = self.tempdir()
        with patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True):
            result = write_generated_file(
                "decks/overview",
                "# Slide One\n\n- Point A\n- Point B\n\n## Slide Two\n\nBody text",
                file_type="pptx",
            )

        expected = Path(tmp).resolve() / "decks" / "overview.pptx"
        self.assertEqual(result.path, expected)
        self.assertEqual(result.file_type, "pptx")
        presentation = Presentation(str(expected))
        self.assertEqual(len(presentation.slides), 2)
        titles = [slide.shapes.title.text for slide in presentation.slides]
        self.assertEqual(titles, ["Slide One", "Slide Two"])

    def test_write_generated_file_maps_legacy_doc_and_ppt_types(self):
        tmp = self.tempdir()
        with patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True):
            doc_result = write_generated_file("legacy-doc.doc", "# Legacy Doc")
            ppt_result = write_generated_file("legacy-deck", "# Legacy Deck", file_type="ppt")

        self.assertEqual(doc_result.path, Path(tmp).resolve() / "legacy-doc.docx")
        self.assertEqual(doc_result.file_type, "docx")
        self.assertEqual(ppt_result.path, Path(tmp).resolve() / "legacy-deck.pptx")
        self.assertEqual(ppt_result.file_type, "pptx")

    def test_write_generated_file_infers_type_from_filename_extension(self):
        tmp = self.tempdir()
        with patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True):
            pdf_result = write_generated_file("reports/inferred.pdf", "PDF body")
            txt_result = write_generated_file("reports/inferred.txt", "text body")
            docx_result = write_generated_file("reports/inferred.docx", "# Word body")

        self.assertEqual(pdf_result.file_type, "pdf")
        self.assertEqual(txt_result.file_type, "txt")
        self.assertEqual(docx_result.file_type, "docx")

    def test_write_generated_file_requires_env_download_path(self):
        with patch.dict("os.environ", {}, clear=True):
            with self.assertRaisesRegex(ValueError, "Download path not defined"):
                write_generated_file("notes", "content")

    def test_file_output_env_takes_precedence_over_download_env(self):
        file_output = self.tempdir()
        download = self.tempdir()
        with patch.dict(
            "os.environ",
            {
                "LOCAL_MCP_FILE_OUTPUT_DIR": file_output,
                "LOCAL_MCP_DOWNLOAD_DIR": download,
            },
            clear=True,
        ):
            result = write_generated_file("preferred-env", "content")

        self.assertEqual(result.path, Path(file_output).resolve() / "preferred-env.md")

    def test_append_generated_file_adds_chunk_to_existing_markdown(self):
        tmp = self.tempdir()
        with patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True):
            write_generated_file("chunks", "first")
            result = append_generated_file("chunks", "second")

        expected = Path(tmp).resolve() / "chunks.md"
        self.assertEqual(expected.read_text(encoding="utf-8"), "first\nsecond\n")
        self.assertEqual(result.path, expected)
        self.assertEqual(result.operation, "append")

    def test_append_generated_file_rejects_binary_output(self):
        tmp = self.tempdir()
        with patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True):
            with self.assertRaisesRegex(ValueError, "Appending to PDF output is not supported"):
                append_generated_file("chunks.pdf", "second")
            with self.assertRaisesRegex(ValueError, "Appending to DOCX output is not supported"):
                append_generated_file("chunks.docx", "second")
            with self.assertRaisesRegex(ValueError, "Appending to PPTX output is not supported"):
                append_generated_file("chunks.pptx", "second")


class GenerateFileToolTests(unittest.IsolatedAsyncioTestCase):
    def tempdir(self):
        root = Path.cwd() / ".tmp" / "unit-tests" / f"{self._testMethodName}-{uuid.uuid4().hex}"
        root.mkdir(parents=True)
        return str(root)

    async def test_generate_file_requires_exactly_one_of_content_and_query(self):
        tmp = self.tempdir()
        with patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True):
            with self.assertRaisesRegex(Exception, "exactly one of `content` or `query`"):
                await generate_file("report")
            with self.assertRaisesRegex(Exception, "exactly one of `content` or `query`"):
                await generate_file("report", content="body", query="a question")

    async def test_generate_file_writes_supplied_content(self):
        tmp = self.tempdir()
        with patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True):
            response = await generate_file("report", content="# Report")

        expected = Path(tmp).resolve() / "report.md"
        self.assertEqual(expected.read_text(encoding="utf-8"), "# Report\n")
        self.assertIn("Markdown file generated.", response)

    async def test_generate_file_append_mode_supports_chunked_writes(self):
        tmp = self.tempdir()
        with patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True):
            await generate_file("report", content="# Report", write_mode="write")
            response = await generate_file("report", content="Second chunk", write_mode="append")

        expected = Path(tmp).resolve() / "report.md"
        self.assertEqual(expected.read_text(encoding="utf-8"), "# Report\nSecond chunk\n")
        self.assertIn("Write mode: append", response)

    async def test_generate_file_supports_pdf_output(self):
        tmp = self.tempdir()
        with patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True):
            response = await generate_file("brief.pdf", content="# Brief\n\nGenerated body")

        expected = Path(tmp).resolve() / "brief.pdf"
        self.assertTrue(expected.is_file())
        self.assertIn("PDF file generated.", response)
        self.assertIn("File type: pdf", response)

    async def test_generate_file_supports_word_and_powerpoint_output(self):
        tmp = self.tempdir()
        with patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True):
            word_response = await generate_file("brief.docx", content="# Brief\n\nWord body")
            deck_response = await generate_file("deck", content="# Deck\n\n- Point", file_type="pptx")

        self.assertTrue((Path(tmp).resolve() / "brief.docx").is_file())
        self.assertTrue((Path(tmp).resolve() / "deck.pptx").is_file())
        self.assertIn("Word file generated.", word_response)
        self.assertIn("File type: docx", word_response)
        self.assertIn("PowerPoint file generated.", deck_response)
        self.assertIn("File type: pptx", deck_response)

    async def test_generate_file_rejects_short_content_when_min_words_is_set(self):
        tmp = self.tempdir()
        with patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True):
            with self.assertRaisesRegex(Exception, "Content is too short"):
                await generate_file("brief.pdf", content="# Brief\n\nToo short.", min_words=50)

        self.assertFalse((Path(tmp).resolve() / "brief.pdf").exists())

    async def test_generate_file_query_mode_uses_smart_search_by_default(self):
        tmp = self.tempdir()
        calls = {}

        async def fake_smart_search(query, **kwargs):
            calls["query"] = query
            calls["kwargs"] = kwargs
            return "Cited summary body. [1]\n\nSources:\n[1] https://example.com/page"

        with (
            patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True),
            patch("local_mcp.tools.file_generation.smart_search.smart_search", new=fake_smart_search),
        ):
            response = await generate_file("research/answer", query="test topic", max_sources=5)

        expected = Path(tmp).resolve() / "research" / "answer.md"
        content = expected.read_text(encoding="utf-8")
        self.assertEqual(calls["query"], "test topic")
        self.assertEqual(calls["kwargs"].get("max_sources"), 5)
        self.assertIn("# test topic", content)
        self.assertIn("Cited summary body.", content)
        self.assertIn("Query: test topic", response)
        self.assertIn("Search mode: smart", response)

    async def test_generate_file_query_mode_supports_deep_research(self):
        tmp = self.tempdir()
        calls = {}

        async def fake_deep_research(query, **kwargs):
            calls["query"] = query
            calls["kwargs"] = kwargs
            return "# Deep Report\n\nFindings body. [1]\n\nSources:\n[1] https://example.com/source"

        with (
            patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True),
            patch("local_mcp.tools.file_generation.deep_research.deep_research", new=fake_deep_research),
        ):
            response = await generate_file(
                "research/deep-report.pdf",
                query="deep topic",
                search_mode="deep",
            )

        expected = Path(tmp).resolve() / "research" / "deep-report.pdf"
        self.assertTrue(expected.is_file())
        self.assertEqual(calls["query"], "deep topic")
        text = "\n".join(page.extract_text() or "" for page in PdfReader(str(expected)).pages)
        self.assertIn("Deep Report", text)
        self.assertIn("Search mode: deep", response)

    async def test_generate_file_query_mode_writes_word_document(self):
        tmp = self.tempdir()

        async def fake_smart_search(query, **kwargs):
            return "Research answer body."

        with (
            patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True),
            patch("local_mcp.tools.file_generation.smart_search.smart_search", new=fake_smart_search),
        ):
            await generate_file("research/word-answer", query="word topic", file_type="doc")

        expected = Path(tmp).resolve() / "research" / "word-answer.docx"
        texts = [paragraph.text for paragraph in Document(str(expected)).paragraphs]
        self.assertIn("word topic", texts)
        self.assertIn("Research answer body.", texts)

    async def test_generate_file_rejects_unknown_search_mode(self):
        tmp = self.tempdir()
        with patch.dict("os.environ", {"LOCAL_MCP_FILE_OUTPUT_DIR": tmp, "LOCAL_MCP_DOWNLOAD_DIR": ""}, clear=True):
            with self.assertRaisesRegex(Exception, "search_mode must be 'smart' or 'deep'"):
                await generate_file("report", query="topic", search_mode="turbo")


if __name__ == "__main__":
    unittest.main()
